"""Document converter — transforms uploaded files into Markdown for LLM consumption.

Supports PDF, Excel, CSV, Word, PowerPoint, HTML, and plain text.
Optional dependencies (pymupdf, openpyxl, python-docx, python-pptx, markdownify)
are imported lazily with graceful fallback errors when missing.

Size limits:
- Maximum file size: 10 MB
- Maximum PDF pages: 50
- Maximum Excel/CSV rows: 10,000
"""

from __future__ import annotations

import csv
import io
import os
from dataclasses import dataclass, field
from pathlib import Path


# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

MAX_FILE_SIZE_BYTES = 10 * 1024 * 1024  # 10 MB
MAX_PDF_PAGES = 50
MAX_SPREADSHEET_ROWS = 10_000


# ---------------------------------------------------------------------------
# Exceptions
# ---------------------------------------------------------------------------


class DocumentConversionError(Exception):
    """Base exception for document conversion failures."""


class UnsupportedFormatError(DocumentConversionError):
    """Raised when the file extension is not in SUPPORTED_TYPES."""


class FileTooLargeError(DocumentConversionError):
    """Raised when the file exceeds MAX_FILE_SIZE_BYTES."""


class DependencyMissingError(DocumentConversionError):
    """Raised when an optional dependency required for conversion is not installed."""


class ContentLimitError(DocumentConversionError):
    """Raised when the document exceeds content limits (pages, rows)."""


# ---------------------------------------------------------------------------
# Result dataclass
# ---------------------------------------------------------------------------


@dataclass
class ConvertedDocument:
    """Result of a successful document conversion."""

    original_path: str
    markdown_path: str
    markdown_content: str
    file_type: str
    page_count: int | None = None
    row_count: int | None = None
    metadata: dict = field(default_factory=dict)


# ---------------------------------------------------------------------------
# Converter
# ---------------------------------------------------------------------------


class DocumentConverter:
    """Convert documents to Markdown for LLM consumption.

    Usage::

        converter = DocumentConverter()
        result = await converter.convert("/path/to/file.csv")
        print(result.markdown_content)
    """

    SUPPORTED_TYPES: dict[str, str] = {
        ".pdf": "_convert_pdf",
        ".xlsx": "_convert_excel",
        ".xls": "_convert_excel",
        ".csv": "_convert_csv",
        ".docx": "_convert_word",
        ".pptx": "_convert_powerpoint",
        ".html": "_convert_html",
        ".htm": "_convert_html",
        ".txt": "_convert_text",
    }

    def __init__(self, output_dir: str | None = None) -> None:
        """Initialise the converter.

        Args:
            output_dir: Directory where converted .md files are written.
                        Defaults to the same directory as the source file.
        """
        self._output_dir = output_dir

    # ------------------------------------------------------------------
    # Public API
    # ------------------------------------------------------------------

    async def convert(self, file_path: str) -> ConvertedDocument:
        """Convert a file to Markdown.

        Args:
            file_path: Absolute or relative path to the source document.

        Returns:
            ConvertedDocument with the markdown content and metadata.

        Raises:
            UnsupportedFormatError: Extension not supported.
            FileTooLargeError: File exceeds 10 MB.
            DependencyMissingError: Required package not installed.
            ContentLimitError: Document exceeds page/row limits.
        """
        path = Path(file_path)
        ext = path.suffix.lower()

        # Validate extension
        if ext not in self.SUPPORTED_TYPES:
            raise UnsupportedFormatError(
                f"Unsupported file format '{ext}'. "
                f"Supported: {', '.join(sorted(self.SUPPORTED_TYPES))}"
            )

        # Validate file exists
        if not path.exists():
            raise DocumentConversionError(f"File not found: {file_path}")

        # Validate size
        file_size = path.stat().st_size
        if file_size > MAX_FILE_SIZE_BYTES:
            size_mb = file_size / (1024 * 1024)
            raise FileTooLargeError(
                f"File size {size_mb:.1f} MB exceeds maximum of "
                f"{MAX_FILE_SIZE_BYTES / (1024 * 1024):.0f} MB"
            )

        # Dispatch to converter method
        method_name = self.SUPPORTED_TYPES[ext]
        method = getattr(self, method_name)
        return await method(path)

    async def convert_bytes(
        self, data: bytes, filename: str, save_dir: str | None = None
    ) -> ConvertedDocument:
        """Convert in-memory bytes to Markdown.

        Writes bytes to a temp file then delegates to convert().

        Args:
            data: Raw file bytes.
            filename: Original filename (used for extension detection).
            save_dir: Directory to save the temp file. Defaults to /tmp.

        Returns:
            ConvertedDocument with the markdown content.
        """
        if len(data) > MAX_FILE_SIZE_BYTES:
            size_mb = len(data) / (1024 * 1024)
            raise FileTooLargeError(
                f"File size {size_mb:.1f} MB exceeds maximum of "
                f"{MAX_FILE_SIZE_BYTES / (1024 * 1024):.0f} MB"
            )

        target_dir = save_dir or self._output_dir or "/tmp"
        os.makedirs(target_dir, exist_ok=True)
        target_path = Path(target_dir) / filename
        target_path.write_bytes(data)
        try:
            return await self.convert(str(target_path))
        finally:
            # Clean up temp file after conversion
            if target_path.exists():
                target_path.unlink()

    # ------------------------------------------------------------------
    # Private converters
    # ------------------------------------------------------------------

    def _md_output_path(self, source: Path) -> Path:
        """Compute the output .md path."""
        out_dir = Path(self._output_dir) if self._output_dir else source.parent
        return out_dir / (source.stem + ".md")

    async def _convert_pdf(self, path: Path) -> ConvertedDocument:
        try:
            import fitz  # pymupdf
        except ImportError:
            raise DependencyMissingError(
                "PDF conversion requires 'pymupdf'. "
                "Install it with: pip install 'agent-orchestrator[docs]'"
            )

        doc = fitz.open(str(path))
        page_count = len(doc)

        if page_count > MAX_PDF_PAGES:
            doc.close()
            raise ContentLimitError(f"PDF has {page_count} pages, maximum is {MAX_PDF_PAGES}")

        parts: list[str] = []
        for i, page in enumerate(doc):
            text = page.get_text()
            if text.strip():
                parts.append(f"## Page {i + 1}\n\n{text.strip()}")
        doc.close()

        md_content = "\n\n".join(parts) if parts else "*Empty PDF document*"
        md_path = self._md_output_path(path)
        md_path.write_text(md_content, encoding="utf-8")

        return ConvertedDocument(
            original_path=str(path),
            markdown_path=str(md_path),
            markdown_content=md_content,
            file_type="pdf",
            page_count=page_count,
        )

    async def _convert_excel(self, path: Path) -> ConvertedDocument:
        try:
            import openpyxl
        except ImportError:
            raise DependencyMissingError(
                "Excel conversion requires 'openpyxl'. "
                "Install it with: pip install 'agent-orchestrator[docs]'"
            )

        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        parts: list[str] = []
        total_rows = 0

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue

            total_rows += len(rows)
            if total_rows > MAX_SPREADSHEET_ROWS:
                wb.close()
                raise ContentLimitError(
                    f"Spreadsheet has {total_rows}+ rows, maximum is {MAX_SPREADSHEET_ROWS}"
                )

            parts.append(f"## Sheet: {sheet_name}\n")
            parts.append(self._rows_to_md_table(rows))

        wb.close()

        md_content = "\n\n".join(parts) if parts else "*Empty spreadsheet*"
        md_path = self._md_output_path(path)
        md_path.write_text(md_content, encoding="utf-8")

        return ConvertedDocument(
            original_path=str(path),
            markdown_path=str(md_path),
            markdown_content=md_content,
            file_type="excel",
            row_count=total_rows,
        )

    async def _convert_csv(self, path: Path) -> ConvertedDocument:
        text = path.read_text(encoding="utf-8", errors="replace")
        reader = csv.reader(io.StringIO(text))
        rows = list(reader)

        if len(rows) > MAX_SPREADSHEET_ROWS:
            raise ContentLimitError(f"CSV has {len(rows)} rows, maximum is {MAX_SPREADSHEET_ROWS}")

        md_content = self._rows_to_md_table(rows) if rows else "*Empty CSV file*"
        md_path = self._md_output_path(path)
        md_path.write_text(md_content, encoding="utf-8")

        return ConvertedDocument(
            original_path=str(path),
            markdown_path=str(md_path),
            markdown_content=md_content,
            file_type="csv",
            row_count=len(rows),
        )

    async def _convert_word(self, path: Path) -> ConvertedDocument:
        try:
            import docx
        except ImportError:
            raise DependencyMissingError(
                "Word conversion requires 'python-docx'. "
                "Install it with: pip install 'agent-orchestrator[docs]'"
            )

        doc = docx.Document(str(path))
        parts: list[str] = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Map Word heading styles to Markdown headings
            style = para.style.name if para.style else ""
            if style.startswith("Heading 1"):
                parts.append(f"# {text}")
            elif style.startswith("Heading 2"):
                parts.append(f"## {text}")
            elif style.startswith("Heading 3"):
                parts.append(f"### {text}")
            else:
                parts.append(text)

        md_content = "\n\n".join(parts) if parts else "*Empty Word document*"
        md_path = self._md_output_path(path)
        md_path.write_text(md_content, encoding="utf-8")

        return ConvertedDocument(
            original_path=str(path),
            markdown_path=str(md_path),
            markdown_content=md_content,
            file_type="docx",
        )

    async def _convert_powerpoint(self, path: Path) -> ConvertedDocument:
        try:
            from pptx import Presentation
        except ImportError:
            raise DependencyMissingError(
                "PowerPoint conversion requires 'python-pptx'. "
                "Install it with: pip install 'agent-orchestrator[docs]'"
            )

        prs = Presentation(str(path))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
            if texts:
                parts.append(f"## Slide {i}\n\n" + "\n\n".join(texts))

        md_content = "\n\n".join(parts) if parts else "*Empty PowerPoint presentation*"
        md_path = self._md_output_path(path)
        md_path.write_text(md_content, encoding="utf-8")

        return ConvertedDocument(
            original_path=str(path),
            markdown_path=str(md_path),
            markdown_content=md_content,
            file_type="pptx",
            page_count=len(prs.slides),
        )

    async def _convert_html(self, path: Path) -> ConvertedDocument:
        text = path.read_text(encoding="utf-8", errors="replace")
        try:
            from markdownify import markdownify
        except ImportError:
            raise DependencyMissingError(
                "HTML conversion requires 'markdownify'. "
                "Install it with: pip install 'agent-orchestrator[docs]'"
            )

        md_content = markdownify(text).strip()
        if not md_content:
            md_content = "*Empty HTML document*"

        md_path = self._md_output_path(path)
        md_path.write_text(md_content, encoding="utf-8")

        return ConvertedDocument(
            original_path=str(path),
            markdown_path=str(md_path),
            markdown_content=md_content,
            file_type="html",
        )

    async def _convert_text(self, path: Path) -> ConvertedDocument:
        text = path.read_text(encoding="utf-8", errors="replace").strip()
        md_content = text if text else "*Empty text file*"

        md_path = self._md_output_path(path)
        md_path.write_text(md_content, encoding="utf-8")

        return ConvertedDocument(
            original_path=str(path),
            markdown_path=str(md_path),
            markdown_content=md_content,
            file_type="txt",
        )

    # ------------------------------------------------------------------
    # Helpers
    # ------------------------------------------------------------------

    @staticmethod
    def _rows_to_md_table(rows: list[tuple | list]) -> str:
        """Convert a list of rows to a Markdown table string."""
        if not rows:
            return ""

        # First row as header
        header = [str(c) if c is not None else "" for c in rows[0]]
        lines = [
            "| " + " | ".join(header) + " |",
            "| " + " | ".join("---" for _ in header) + " |",
        ]
        for row in rows[1:]:
            cells = [str(c) if c is not None else "" for c in row]
            # Pad or truncate to match header width
            while len(cells) < len(header):
                cells.append("")
            cells = cells[: len(header)]
            lines.append("| " + " | ".join(cells) + " |")

        return "\n".join(lines)
